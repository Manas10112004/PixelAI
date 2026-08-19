from pptx import Presentation

def create_hackathon_ppt():
    prs = Presentation()

    # Slide 1: Title
    slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide_1.shapes.title.text = "PixelAI: AI-Based Restoration of Degraded Images"
    slide_1.placeholders[1].text = "HPC-OPTIMIZED • NAFNET • IMAGE RESTORATION\nTeam: Gresha Mandape, Manas Bhole, Shradha Joshi, Vinaya Gabale\nIndira University"

    # Slide 2: Problem
    slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_2.shapes.title.text = "01 • THE CHALLENGE"
    tf = slide_2.placeholders[1].text_frame
    tf.text = "Problem Statement Addressed: AI-Based Restoration of Degraded Images"
    tf.add_paragraph().text = "High-resolution wafer scans are essential for micro-defect detection."
    tf.add_paragraph().text = "Sensor noise and downsampling obscure critical circuit details."
    tf.add_paragraph().text = "Goal: fast, high-throughput restoration without blur or hallucinated structure."

    # Slide 3: Concept
    slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_3.shapes.title.text = "02 • CORE CONCEPT"
    tf = slide_3.placeholders[1].text_frame
    tf.text = "Idea Description: SemiconNAFNet + algorithmic emergency-rescue filter"
    tf.add_paragraph().text = "Pipeline: Degraded Scan -> SemiconNAFNet -> PixelShuffle -> Restored"
    tf.add_paragraph().text = "Activation-free design reduces memory round-trips for efficient GPU execution."
    tf.add_paragraph().text = "Targets: Speckle/Gaussian, Impulse noise, and Super-resolution."

    # Slide 4: Architecture
    slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_4.shapes.title.text = "03 • ARCHITECTURE & TRAINING"
    tf = slide_4.placeholders[1].text_frame
    tf.text = "Proposed Solution"
    tf.add_paragraph().text = "Model: Input -> SimpleGate -> Depthwise Blocks -> PixelShuffle -> Output"
    tf.add_paragraph().text = "Training: torch.amp.autocast + Cosine Annealing LR"
    tf.add_paragraph().text = "Loss: Direct L1 + numerically safe 2D FFT frequency loss"

    # Slide 5: Innovation
    slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_5.shapes.title.text = "04 • WHAT DIFFERENTIATES US"
    tf = slide_5.placeholders[1].text_frame
    tf.text = "Innovation & Uniqueness"
    tf.add_paragraph().text = "HPC-first design: Tensor-Core aligned channels + memory-mapped ingestion"
    tf.add_paragraph().text = "Safe spectral loss avoids unstable FFT magnitudes"
    tf.add_paragraph().text = "Rescue bypass: OpenCV fallback (fastNlMeansDenoising + medianBlur)"

    # Slide 6: Results
    slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_6.shapes.title.text = "05 • EVIDENCE"
    tf = slide_6.placeholders[1].text_frame
    tf.text = "Quantitative Results"
    tf.add_paragraph().text = "PSNR: 28.24 dB"
    tf.add_paragraph().text = "SSIM: 0.7222"
    tf.add_paragraph().text = "Visual Evidence: Complete noise-spike reduction with preserved structural integrity."

    # Slide 7: Implementation
    slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_7.shapes.title.text = "06 • IMPLEMENTATION"
    tf = slide_7.placeholders[1].text_frame
    tf.text = "Technology & Feasibility"
    tf.add_paragraph().text = "Tech stack: PyTorch, OpenCV, NumPy, PIL, Matplotlib"
    tf.add_paragraph().text = "Hardware: NVIDIA GPU (CUDA-accelerated)"
    tf.add_paragraph().text = "Model: Lightweight activation-free design"
    tf.add_paragraph().text = "Inference: Milliseconds per tile"

    # Slide 8: Submission
    slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_8.shapes.title.text = "07 • SUBMISSION"
    tf = slide_8.placeholders[1].text_frame
    tf.text = "Links"
    tf.add_paragraph().text = "GitHub: https://github.com/Manas10112004/PixelAI"
    tf.add_paragraph().text = "Demo Video: [Link]"

    # Slide 9: Sources
    slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_9.shapes.title.text = "08 • SOURCES"
    tf = slide_9.placeholders[1].text_frame
    tf.text = "References"
    tf.add_paragraph().text = "Chen, L. et al. — Simple Baselines for Image Restoration"
    tf.add_paragraph().text = "PyTorch Documentation — Mixed Precision Training"
    tf.add_paragraph().text = "OpenCV Documentation — Non-Local Means Denoising"

    prs.save("PixelAI_Final_Submission.pptx")
    print("[✓] Presentation saved as PixelAI_Final_Submission.pptx")

if __name__ == "__main__":
    create_hackathon_ppt()